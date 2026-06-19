# backend/utils.py
import yfinance as yf
# Import yahooquery lazily inside resolve_ticker_dynamically to avoid import-time failures
# pptx imports are done lazily inside add_institutional_footer to avoid import-time failure when python-pptx isn't installed

def resolve_ticker_dynamically(ticker):
    """Fixes common ticker errors (e.g. AIRBUS -> AIR.PA)."""
    try:
        t = yf.Ticker(ticker)
        if t.info and 'regularMarketPrice' in t.info: return ticker
    except: pass
    try:
        try:
            from yahooquery import search
        except Exception:
            search = None
        if search:
            data = search(ticker)
            if data and 'quotes' in data and len(data['quotes']) > 0:
                return data['quotes'][0]['symbol']
    except: pass
    return ticker

def get_currency_symbol(currency_code):
    if currency_code == 'GBp': return '£'
    symbols = { 'USD': '$', 'INR': '₹', 'EUR': '€', 'GBP': '£', 'JPY': '¥', 'CAD': 'C$', 'AUD': 'A$', 'CNY': '¥' }
    return symbols.get(currency_code, currency_code + " ")

def format_large_number(num):
    if not num or isinstance(num, str): return "N/A"
    try:
        val = float(num)
        if val > 1e12: return f"{round(val/1e12, 2)}T"
        elif val > 1e9: return f"{round(val/1e9, 2)}B"
        else: return f"{round(val/1e6, 2)}M"
    except: return "N/A"

def get_exchange_rate(from_curr, to_curr):
    """Fetches live FX rates, handling Pence (GBp)."""
    if not from_curr or not to_curr or from_curr == to_curr: return 1.0
    if from_curr == 'GBp' and to_curr == 'GBP': return 0.01
    if from_curr == 'GBP' and to_curr == 'GBp': return 100.0
    if from_curr == 'GBp': return 0.01 * get_exchange_rate('GBP', to_curr)
    
    try:
        pair = f"{from_curr}{to_curr}=X"
        hist = yf.Ticker(pair).history(period="1d")
        if not hist.empty: return hist['Close'].iloc[-1]
        
        pair_inv = f"{to_curr}{from_curr}=X"
        hist = yf.Ticker(pair_inv).history(period="1d")
        if not hist.empty: return 1.0 / hist['Close'].iloc[-1]
    except: pass
    return 1.0

def normalize_data(info):
    """Normalizes price and currency (Fixes Pence/Pounds issue)."""
    price = info.get('currentPrice') or info.get('regularMarketPrice')
    currency = info.get('currency', 'USD')
    if currency == 'GBp':
        price = price / 100.0 if price else 0
        currency = 'GBP'
    return price, currency

def add_institutional_footer(slide):
    """Adds a small footer. Import pptx utilities lazily so missing library doesn't break app import."""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception:
        # If python-pptx is not available, try to add an unformatted textbox, but silently return if impossible
        try:
            txBox = slide.shapes.add_textbox(0, 0, 0, 0)
            txBox.text_frame.paragraphs[0].text = "CONFIDENTIAL | FOR EDUCATIONAL PURPOSES ONLY | SAMPADA"
        except Exception:
            pass
        return

    left = Inches(0.5); top = Inches(7.2); width = Inches(9.0); height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = "CONFIDENTIAL | FOR EDUCATIONAL PURPOSES ONLY | SAMPADA"
    p.font.size = Pt(8); p.font.color.rgb = RGBColor(100, 100, 100); p.alignment = PP_ALIGN.CENTER