# backend/main.py
import os
import re
import time
import math
from io import BytesIO
from typing import Optional

import yfinance as yf
from fastapi import FastAPI, HTTPException, Header, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.exception_handlers import http_exception_handler
from pydantic import BaseModel, validator

from config import CACHE, CACHE_DURATION, LEGAL_DISCLAIMER
from models import ExportRequest
from utils import (
    resolve_ticker_dynamically, normalize_data, get_currency_symbol,
    format_large_number, add_institutional_footer, get_exchange_rate,
)
from services import AlphaVantageService
from analysis import (
    run_institutional_dcf, get_robust_peers, process_peer_data,
    calculate_weighted_harmonic_mean, get_smart_news, get_financial_statements,
    get_historical_trading_range, get_analyst_consensus,
    calculate_quality_scores, generate_swot, calculate_confidence_score,
    calculate_blend_weights, run_backtest,
    _safe_float,
)
from auth import run_startup_diagnostic, verify_token, save_session

# ---------------------------------------------------------------------------
# SECURITY: allowed CORS origins from env (defaults to localhost for dev)
# In production set ALLOWED_ORIGINS=https://yourdomain.com,https://app.yourdomain.com
# ---------------------------------------------------------------------------
_raw_origins = os.environ.get("ALLOWED_ORIGINS", "http://localhost:3000,http://127.0.0.1:3000")
ALLOWED_ORIGINS = [o.strip() for o in _raw_origins.split(",") if o.strip()]

# ---------------------------------------------------------------------------
# SECURITY: strict ticker validation — only alphanumeric + . and - (e.g. AIR.PA, BRK-B)
# ---------------------------------------------------------------------------
_TICKER_RE = re.compile(r'^[A-Za-z0-9]([A-Za-z0-9.\-]{0,14}[A-Za-z0-9])?$')


def _validate_ticker(raw: str) -> str:
    """Sanitise and validate a ticker symbol. Raises HTTPException 400 on invalid input."""
    clean = raw.strip().upper()[:20]
    if not _TICKER_RE.match(clean):
        raise HTTPException(status_code=400, detail="Invalid ticker symbol format.")
    return clean


# ---------------------------------------------------------------------------
# APP INIT & STARTUP
# ---------------------------------------------------------------------------
app = FastAPI(title="SAMPADA.ai API", version="2.1.0", docs_url=None, redoc_url=None)

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type"],
)


# ---------------------------------------------------------------------------
# SECURITY: suppress internal stack traces from HTTP error responses
# ---------------------------------------------------------------------------
@app.exception_handler(Exception)
async def _generic_exception_handler(request: Request, exc: Exception):
    # Log internally but never leak traceback to client
    print(f"[ERROR] {request.url.path}: {type(exc).__name__}: {exc}")
    return JSONResponse(status_code=500, content={"detail": "An internal error occurred."})


@app.exception_handler(HTTPException)
async def _http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse(status_code=exc.status_code, content={"detail": exc.detail})

_firebase_diag: dict = {}


@app.on_event("startup")
def _startup():
    global _firebase_diag
    print("[SAMPADA] Starting up …")
    _firebase_diag = run_startup_diagnostic()
    print(f"[SAMPADA] Firebase status: {_firebase_diag.get('status')}")


# ---------------------------------------------------------------------------
# DIAGNOSTIC ENDPOINT
# ---------------------------------------------------------------------------

@app.get("/")
def health_check():
    return {"status": "ok", "app": "SAMPADA.ai"}


@app.get("/api/diagnostics")
def diagnostics():
    return {
        "app": "SAMPADA.ai",
        "version": "2.1.0",
        "firebase": _firebase_diag,
        "cache_entries": len(CACHE),
    }


# ---------------------------------------------------------------------------
# FIREBASE AUTH ENDPOINTS
# ---------------------------------------------------------------------------

class TokenBody(BaseModel):
    id_token: str


@app.post("/api/auth/verify")
def auth_verify(body: TokenBody):
    try:
        claims = verify_token(body.id_token)
        return {"uid": claims.get("uid"), "email": claims.get("email"), "valid": True}
    except Exception as exc:
        raise HTTPException(status_code=401, detail=str(exc))


# ---------------------------------------------------------------------------
# MAIN ANALYSIS ENDPOINT
# ---------------------------------------------------------------------------

def _generate_theses(info, dcf_upside, sentiment_data):
    rev_g = _safe_float(info.get('revenueGrowth'), 0) or 0
    margin = _safe_float(info.get('profitMargins'), 0) or 0
    debt_eq = _safe_float(info.get('debtToEquity'), 0) or 0
    pe = _safe_float(info.get('trailingPE'), 0) or 0

    bull = []
    if rev_g > 0.10:
        bull.append(f"Strong double-digit revenue growth ({rev_g:.1%}) signals expanding market share.")
    elif rev_g > 0.05:
        bull.append("Steady top-line expansion supports stable cash flow generation.")
    if margin > 0.15:
        bull.append(f"High profit margins ({margin:.1%}) demonstrate strong pricing power.")
    if dcf_upside > 15:
        bull.append("Valuation is attractive with >15% upside to intrinsic DCF value.")
    if (sentiment_data.get('short_term') or 50) > 60:
        bull.append("Short-term sentiment momentum is positive, supported by recent news flow.")
    if not bull:
        bull.append("Defensive positioning in a volatile market environment.")

    bear = []
    if debt_eq > 150:
        bear.append(f"Elevated leverage (D/E: {debt_eq:.0f}%) increases sensitivity to rate hikes.")
    if pe > 40:
        bear.append(f"Rich valuation (P/E: {pe:.1f}x) leaves little room for execution error.")
    elif pe > 25:
        bear.append("Trading at a premium multiple relative to historical averages.")
    if margin < 0.05:
        bear.append(f"Thin margins ({margin:.1%}) make earnings vulnerable to cost inflation.")
    if sentiment_data.get('event_risk') == 'High':
        bear.append("High event risk detected in recent news (Regulatory/Legal/Macro).")
    if not bear:
        bear.append("Broader sector rotation risks could limit near-term multiple expansion.")

    return {"bull": " ".join(bull), "bear": " ".join(bear)}


@app.get("/api/analyze/{ticker}")
def analyze_stock(ticker: str, authorization: Optional[str] = Header(default=None)):
    ticker = _validate_ticker(ticker)
    ticker = resolve_ticker_dynamically(ticker)

    cached = CACHE.get(ticker)
    if cached and (time.time() - cached['timestamp']) < CACHE_DURATION:
        return cached['data']

    is_us = "." not in ticker

    try:
        stock = yf.Ticker(ticker)
        info = stock.info
        price_check = info.get('regularMarketPrice') or info.get('currentPrice')
        if not info or not price_check:
            raise ValueError("No market data")
    except Exception:
        raise HTTPException(status_code=404, detail="Ticker not found or no market data available")

    price, currency = normalize_data(info)
    mkt_cap = info.get('marketCap')

    macro = AlphaVantageService.get_macro_indicators()
    commodities = AlphaVantageService.get_commodities()

    SECTOR_PE_MAP = {
        "Technology": 25, "Financial Services": 12, "Healthcare": 20,
        "Consumer Cyclical": 18, "Industrials": 18, "Energy": 14,
        "Real Estate": 22, "Utilities": 18, "Communication Services": 20,
        "Consumer Defensive": 20, "Basic Materials": 16,
    }
    sector_pe = SECTOR_PE_MAP.get(info.get('sector'), 20)

    dcf_data, dcf_method = run_institutional_dcf(ticker, info, stock, sector_pe)

    target_prof = {
        "mkt_cap": mkt_cap,
        "growth": _safe_float(info.get('revenueGrowth'), 0),
        "margin": _safe_float(info.get('ebitdaMargins'), 0),
        "roic": _safe_float(info.get('returnOnEquity'), 0),
        "capex_intensity": abs(_safe_float(info.get('capitalExpenditures'), 0) or 0) / (
            (_safe_float(info.get('totalRevenue'), 1) or 1)),
        "net_debt_ebitda": (
            (_safe_float(info.get('totalDebt'), 0) or 0) -
            (_safe_float(info.get('totalCash'), 0) or 0)
        ) / ((_safe_float(info.get('ebitda'), 1) or 1)),
    }

    # ── Peer Discovery (Category A + B) ──────────────────────────────────────
    peer_result = get_robust_peers(
        ticker, info.get('sector'), info.get('industry'), mkt_cap, is_us,
        target_revenue=_safe_float(info.get('totalRevenue')),
        target_roic=_safe_float(info.get('returnOnEquity')),
        target_margin=_safe_float(info.get('ebitdaMargins')),
    )
    peer_methodology = peer_result.get("methodology", {})

    all_processed = []
    for cat, tickers_list in [("A", peer_result.get("cat_a", [])),
                               ("B", peer_result.get("cat_b", []))]:
        for p_ticker in tickers_list:
            pd = process_peer_data(p_ticker, currency, target_prof, category=cat)
            if pd:
                all_processed.append(pd)

    all_processed.sort(key=lambda x: x['similarity'], reverse=True)
    # Prioritise Cat A in final display (max 4 Cat A + 2 Cat B)
    cat_a_peers = [p for p in all_processed if p.get('category') == 'A'][:4]
    cat_b_peers = [p for p in all_processed if p.get('category') == 'B'][:2]
    final_peers = cat_a_peers + cat_b_peers

    # ── Valuation Calculation ─────────────────────────────────────────────────
    pe_vals, ev_vals, w_vals = [], [], []
    for pd in final_peers:
        if pd.get('pe'):
            pe_vals.append(pd['pe'])
            w_vals.append(pd['similarity'])
        if pd.get('ev_ebitda'):
            ev_vals.append(pd['ev_ebitda'])

    avg_pe = calculate_weighted_harmonic_mean(pe_vals, w_vals)
    avg_ev = calculate_weighted_harmonic_mean(ev_vals, w_vals)

    comps_price = price or 0
    if avg_ev > 0 and info.get('enterpriseToEbitda'):
        comps_price = price * (avg_ev / info['enterpriseToEbitda'])
    elif avg_pe > 0 and info.get('trailingEps'):
        comps_price = avg_pe * info['trailingEps']

    # Dynamic, sector-aware DCF/Comps blend (replaces the fixed 60/40).
    w_dcf, w_comps = calculate_blend_weights(info, info.get('sector'))
    if dcf_data['val'] > 0 and comps_price > 0:
        final_val = dcf_data['val'] * w_dcf + comps_price * w_comps
    elif dcf_data['val'] > 0:
        final_val = dcf_data['val']
    else:
        final_val = comps_price
        w_dcf, w_comps = 0.0, 1.0
    upside = round(((final_val - price) / price) * 100, 2) if price else 0

    verdict = "NEUTRAL"
    if upside > 15:
        verdict = "POSITIVE BIAS"
    elif upside < -10:
        verdict = "NEGATIVE BIAS"

    sentiment_data, headlines = get_smart_news(ticker, info, is_us)
    theses = _generate_theses(info, upside, sentiment_data)
    financials = get_financial_statements(stock)
    hist = get_historical_trading_range(stock)
    consensus = get_analyst_consensus(ticker, info, is_us)
    quality_scores = calculate_quality_scores(info, financials)
    swot = generate_swot(info, [])
    backtest = run_backtest(
        stock, final_val, price,
        dcf_low=dcf_data.get('low'), dcf_high=dcf_data.get('high'),
        consensus_target=consensus.get('target_mean'),
    )

    football_field = {
        "fifty_two_week": [hist['low_52wk'], hist['high_52wk']] if hist else [price * 0.8, price * 1.2],
        "analyst_target": [
            consensus['target_low'] or price * 0.9,
            consensus['target_high'] or price * 1.1,
        ],
        "dcf_range": [dcf_data['low'], dcf_data['high']],
        "comps_range": [comps_price * 0.85, comps_price * 1.15],
    }

    result = {
        "symbol": ticker,
        "name": info.get('shortName', ticker),
        "currency_symbol": get_currency_symbol(currency),
        "price": round(price, 2),
        "market_cap": format_large_number(mkt_cap),
        "sentiment_score": sentiment_data['score'],
        "sentiment_analysis": sentiment_data,
        "sentiment_explanation": f"Event Risk: {sentiment_data['event_risk']} | Short-Term: {sentiment_data['short_term']}/100",
        "verdict": verdict,
        "confidence_score": calculate_confidence_score(final_peers, dcf_data['val'], info, sentiment_data['event_risk']),
        "headlines": headlines,
        "peers": final_peers,
        "peer_methodology": peer_methodology,
        "summary": info.get('longBusinessSummary'),
        "sector": info.get('sector'),
        "industry": info.get('industry'),
        "website": info.get('website', 'N/A'),
        "employees": format_large_number(info.get('fullTimeEmployees', 0)),
        "city": info.get('city', 'N/A'),
        "historical_data": hist,
        "financials": financials,
        "target_profile": target_prof,
        "target_ratios": {
            "pe": round(info.get('trailingPE', 0), 2) if info.get('trailingPE') else "-",
            "ev_ebitda": round(info.get('enterpriseToEbitda', 0), 2) if info.get('enterpriseToEbitda') else "-",
            "ev_sales": round(info.get('enterpriseToRevenue', 0), 2) if info.get('enterpriseToRevenue') else "-",
            "net_debt_ebitda": round(target_prof['net_debt_ebitda'], 2) if info.get('ebitda') else 0,
            "roic": f"{round((_safe_float(info.get('returnOnEquity')) or 0) * 100, 1)}%",
            "growth": round((_safe_float(info.get('revenueGrowth')) or 0) * 100, 1),
        },
        "valuation_analysis": {
            "implied_price": round(final_val, 2),
            "upside": upside,
            "dcf_price": round(dcf_data['val'], 2),
            "dcf_range": f"{dcf_data['low']} - {dcf_data['high']}",
            "comps_price": round(comps_price, 2),
            "wacc": dcf_data.get('wacc', 0),
            "growth_assumed": dcf_data.get('growth', 0),
            "fcf_lookback": dcf_data.get('fcf_lookback'),
            "blend_dcf": round(w_dcf * 100),
            "blend_comps": round(w_comps * 100),
        },
        "quality_scores": quality_scores,
        "swot": swot,
        "consensus": consensus,
        "backtest": backtest,
        "football_field": football_field,
        "wacc_components": {
            "rf": macro["risk_free_rate"],
            "beta": info.get('beta', 1.0),
            "erp": 0.055,
            "wacc": dcf_data.get('wacc', 0.1),
        },
        "tear_sheet_data": {
            "price": round(price, 2),
            "market_cap": format_large_number(mkt_cap),
            "revenue": format_large_number(info.get('totalRevenue', 0)),
            "ebitda": format_large_number(info.get('ebitda', 0)),
            "gross_profit": format_large_number(info.get('grossProfits', 0)),
            "net_income": format_large_number(info.get('netIncomeToCommon', 0)),
            "fcf": format_large_number(info.get('freeCashflow', 0)),
            "total_debt": format_large_number(info.get('totalDebt', 0)),
            "cash": format_large_number(info.get('totalCash', 0)),
            "employees": format_large_number(info.get('fullTimeEmployees', 0)),
        },
        "market_data": {"commodities": commodities, "macro": macro},
        "theses": theses,
    }

    CACHE[ticker] = {'timestamp': time.time(), 'data': result}

    if authorization and authorization.startswith("Bearer "):
        try:
            claims = verify_token(authorization.split(" ", 1)[1])
            save_session(claims.get("uid", ""), ticker, result)
        except Exception:
            pass

    return result


# ---------------------------------------------------------------------------
# MANUAL PEER LOOKUP
# ---------------------------------------------------------------------------

@app.get("/api/peer_info")
def get_manual_peer(ticker: str, base_currency: str = "USD"):
    ticker = _validate_ticker(ticker)
    res = process_peer_data(ticker, base_currency)
    if not res:
        raise HTTPException(status_code=404, detail="Not Found")
    for k in ("pe", "ev_ebitda", "ev_sales"):
        if res.get(k) and isinstance(res[k], float):
            res[k] = round(res[k], 2)
    return res


# ---------------------------------------------------------------------------
# TICKER SEARCH  (live combobox for peer add)
# ---------------------------------------------------------------------------

_SEARCH_QUERY_RE = re.compile(r'^[A-Za-z0-9 .\-]{1,40}$')


@app.get("/api/ticker_search")
def ticker_search(q: str):
    """
    Keyless live ticker search via yahooquery.search().
    Returns up to 8 results: [{symbol, name, exchange, type}]
    q must be 1–40 chars of alphanumeric / space / dot / dash.
    """
    q = q.strip()
    if not q or not _SEARCH_QUERY_RE.match(q):
        raise HTTPException(status_code=400, detail="Invalid search query.")
    try:
        from yahooquery import search as yq_search
        results = yq_search(q)
        quotes = results.get("quotes", []) if isinstance(results, dict) else []
        filtered = [
            {
                "symbol":   r.get("symbol", ""),
                "name":     r.get("longname") or r.get("shortname") or r.get("symbol", ""),
                "exchange": r.get("exchange", ""),
                "type":     r.get("quoteType", "EQUITY"),
            }
            for r in quotes
            if r.get("symbol") and r.get("quoteType") in ("EQUITY", "ETF", "FUND")
        ][:8]
        return {"results": filtered}
    except Exception:
        return {"results": []}


# ---------------------------------------------------------------------------
# ELITE POWERPOINT EXPORT ENGINE  (Phase 2 — Investment Banking Grade)
# ---------------------------------------------------------------------------

@app.post("/api/export_ppt")
async def export_ppt(request: ExportRequest):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed.")

    is_client = request.view_mode == "Client"
    W, H = 10.0, 7.5  # slide dimensions in inches

    # ── Palette ────────────────────────────────────────────────────────────────
    BG       = RGBColor(0x0B, 0x15, 0x21)   # deep navy
    PANEL    = RGBColor(0x10, 0x1E, 0x2E)   # panel
    PANEL2   = RGBColor(0x16, 0x28, 0x3C)   # lighter panel
    ACCENT   = RGBColor(0x00, 0xD4, 0xFF)   # cyan
    WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT    = RGBColor(0xCC, 0xDD, 0xEE)
    DIM      = RGBColor(0x55, 0x77, 0x99)
    DIMMER   = RGBColor(0x2A, 0x3A, 0x4A)
    GREEN    = RGBColor(0x00, 0xE5, 0x76)
    RED      = RGBColor(0xFF, 0x44, 0x44)
    YELLOW   = RGBColor(0xFF, 0xA5, 0x00)
    HDR_ROW  = RGBColor(0x00, 0x5C, 0x8A)
    HDR_BULL = RGBColor(0x00, 0x5C, 0x2A)
    HDR_BEAR = RGBColor(0x6A, 0x10, 0x10)
    STRIP_A  = RGBColor(0x10, 0x1E, 0x2E)
    STRIP_B  = RGBColor(0x0D, 0x19, 0x28)

    def _verdict_col(v: str) -> RGBColor:
        return GREEN if "POSITIVE" in v else (RED if "NEGATIVE" in v or "REVIEW" in v else YELLOW)

    # ── PPT helpers ────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def _slide() -> object:
        return prs.slides.add_slide(blank)

    def _bg(slide, color: RGBColor):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _rect(slide, lx, ty, w, h, fill: RGBColor, line_color=None):
        from pptx.util import Inches
        s = slide.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if line_color:
            s.line.color.rgb = line_color
            s.line.width = Pt(0.5)
        else:
            s.line.fill.background()
        return s

    def _txt(slide, text: str, lx, ty, w, h,
             size=10, bold=False, color: RGBColor = None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
        from pptx.util import Inches, Pt
        txb = slide.shapes.add_textbox(Inches(lx), Inches(ty), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or WHITE
        return txb

    def _footer(slide):
        _rect(slide, 0, 7.30, W, 0.20, DIMMER)
        _txt(slide,
             f"CONFIDENTIAL  ·  SAMPADA.ai  ·  {time.strftime('%B %d, %Y')}  ·  FOR EDUCATIONAL PURPOSES ONLY",
             0.3, 7.32, W - 0.6, 0.16,
             size=6.5, color=DIM, align=PP_ALIGN.CENTER)

    def _section_header(slide, text: str, ty: float):
        """Thin cyan rule + uppercase label."""
        _rect(slide, 0.35, ty, W - 0.7, 0.03, ACCENT)
        _txt(slide, text, 0.35, ty + 0.06, 8.0, 0.22, size=7.5, bold=True, color=ACCENT)

    def _kv_box(slide, label, value, lx, ty, w=2.05, h=0.72, val_color=None):
        """Styled key-value metric block."""
        _rect(slide, lx, ty, w, h, PANEL2)
        _rect(slide, lx, ty, 0.06, h, ACCENT)
        _txt(slide, label.upper(), lx + 0.12, ty + 0.07, w - 0.18, 0.18,
             size=6.5, color=DIM, bold=True)
        _txt(slide, str(value), lx + 0.12, ty + 0.29, w - 0.18, 0.34,
             size=13, bold=True, color=val_color or WHITE)

    # ── Display verdict ────────────────────────────────────────────────────────
    display_verdict = ("REVIEW" if (is_client and "NEGATIVE" in request.verdict)
                       else request.verdict)
    v_col = _verdict_col(display_verdict)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — COVER
    # ════════════════════════════════════════════════════════════════════════════
    s1 = _slide()
    _bg(s1, BG)

    # Left accent stripe
    _rect(s1, 0, 0, 0.10, H, ACCENT)

    # Header band
    _rect(s1, 0.10, 0, W - 0.10, 1.05, PANEL)
    _txt(s1, "SAMPADA.ai", 0.30, 0.10, 4.0, 0.40, size=11, bold=True, color=ACCENT)
    _txt(s1, f"{'CLIENT PRESENTATION' if is_client else 'INTERNAL — STRICTLY PRIVATE & CONFIDENTIAL'}",
         0.30, 0.55, 7.0, 0.30, size=7.5, color=DIM)
    _txt(s1, time.strftime('%B %d, %Y'), 8.0, 0.55, 1.7, 0.30,
         size=7.5, color=DIM, align=PP_ALIGN.RIGHT)

    # Company name + ticker
    _txt(s1, request.target_name, 0.35, 1.35, 9.0, 1.05,
         size=36, bold=True, color=WHITE, wrap=True)
    _txt(s1, f"[ {request.ticker}  ·  {request.sector or ''}  ·  {request.industry or ''} ]",
         0.35, 2.45, 9.0, 0.35, size=9, color=DIM)

    # Verdict badge
    _rect(s1, 0.35, 2.95, 3.5, 0.85, PANEL2)
    _rect(s1, 0.35, 2.95, 0.12, 0.85, v_col)
    _txt(s1, "ASSESSMENT", 0.55, 2.98, 3.0, 0.22, size=6.5, bold=True, color=DIM)
    _txt(s1, display_verdict, 0.55, 3.19, 3.0, 0.45, size=18, bold=True, color=v_col)

    # Key stats row
    cs = request.currency_symbol
    stats = [
        ("Current Price", f"{cs}{request.current_price:.2f}"),
        ("Implied Value", f"{cs}{request.implied_price:.2f}"),
        ("Spread to Market", f"{request.upside:+.1f}%"),
        ("Deal Type", request.deal_type),
    ]
    stat_w = 2.0
    for i, (lbl, val) in enumerate(stats):
        _kv_box(s1, lbl, val,
                lx=4.30 + i * (stat_w + 0.15), ty=2.95, w=stat_w, h=0.85,
                val_color=GREEN if "+" in str(val) else (RED if "-" in str(val) and i == 2 else WHITE))

    # Summary excerpt
    summ = (request.summary or "")[:280]
    if summ:
        _rect(s1, 0.35, 4.05, 9.3, 1.60, PANEL)
        _txt(s1, "EXECUTIVE SUMMARY", 0.50, 4.10, 8.0, 0.22, size=7, bold=True, color=DIM)
        _txt(s1, summ + ("…" if len(request.summary or "") > 280 else ""),
             0.50, 4.36, 9.0, 1.20, size=8.5, color=LIGHT, wrap=True)

    _footer(s1)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — DISCLAIMER
    # ════════════════════════════════════════════════════════════════════════════
    s2 = _slide()
    _bg(s2, BG)
    _rect(s2, 0, 0, W, 0.85, PANEL)
    _rect(s2, 0, 0, 0.10, H, HDR_BEAR)
    _txt(s2, "IMPORTANT NOTICE & DISCLAIMER", 0.25, 0.22, 9.0, 0.40,
         size=14, bold=True, color=WHITE)
    _txt(s2, LEGAL_DISCLAIMER.strip(), 0.30, 1.05, 9.4, 5.6,
         size=9, color=LIGHT, wrap=True)
    _footer(s2)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — FINANCIAL TEAR SHEET
    # ════════════════════════════════════════════════════════════════════════════
    s3 = _slide()
    _bg(s3, BG)
    _rect(s3, 0, 0, W, 0.70, PANEL)
    _rect(s3, 0, 0, 0.10, H, ACCENT)
    _txt(s3, f"{request.ticker}  ·  Financial Tear Sheet",
         0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

    td = request.tear_sheet_data or {}
    items = [
        ("Current Price",    f"{cs}{td.get('price', '-')}"),
        ("Market Cap",       td.get('market_cap', '-')),
        ("Revenue (LTM)",    td.get('revenue', '-')),
        ("EBITDA (LTM)",     td.get('ebitda', '-')),
        ("Gross Profit",     td.get('gross_profit', '-')),
        ("Net Income",       td.get('net_income', '-')),
        ("Free Cash Flow",   td.get('fcf', '-')),
        ("Total Debt",       td.get('total_debt', '-')),
        ("Cash & Equiv.",    td.get('cash', '-')),
        ("Employees",        td.get('employees', '-')),
    ]
    cols, col_w, row_h = 2, 4.55, 0.62
    for idx, (lbl, val) in enumerate(items):
        col = idx % cols
        row = idx // cols
        lx = 0.25 + col * (col_w + 0.40)
        ty = 0.85 + row * (row_h + 0.06)
        _kv_box(s3, lbl, val, lx=lx, ty=ty, w=col_w, h=row_h)

    _footer(s3)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — VALUATION FOOTBALL FIELD  (visual bars)
    # ════════════════════════════════════════════════════════════════════════════
    s4 = _slide()
    _bg(s4, BG)
    _rect(s4, 0, 0, W, 0.70, PANEL)
    _rect(s4, 0, 0, 0.10, H, ACCENT)
    _txt(s4, f"{request.ticker}  ·  Valuation Football Field  ({request.deal_type})",
         0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

    ff = request.football_field or {}
    cur_price = request.current_price or request.implied_price
    rows = [
        ("52-Week Range",   ff.get("fifty_two_week",  [cur_price * 0.8, cur_price * 1.2]),  PANEL2),
        ("Analyst Targets", ff.get("analyst_target",  [cur_price * 0.9, cur_price * 1.1]),  DIMMER),
        ("Comps Range",     ff.get("comps_range",     [cur_price * 0.85, cur_price * 1.15]), PANEL),
        ("DCF Range",       ff.get("dcf_range",       [cur_price * 0.8, cur_price * 1.3]),  PANEL2),
    ]

    # Global min / max across all ranges
    all_vals = [v for row in rows for v in row[1] if v]
    all_vals.append(cur_price)
    g_min = min(all_vals) * 0.88
    g_max = max(all_vals) * 1.10
    g_range = g_max - g_min or 1

    chart_lx = 1.70   # start of bar area
    chart_w  = 7.80   # width of bar area
    row_start_y = 0.95
    row_h2   = 0.72
    bar_h    = 0.30
    bar_colors = [ACCENT, GREEN, YELLOW, RGBColor(0xAA, 0x55, 0xFF)]

    for ri, (label, rng, bg_col) in enumerate(rows):
        ty = row_start_y + ri * (row_h2 + 0.04)
        lo, hi = (rng[0] or g_min), (rng[1] or g_max)
        lo_pct = (lo - g_min) / g_range
        hi_pct = (hi - g_min) / g_range
        bar_lx = chart_lx + lo_pct * chart_w
        bar_w  = max((hi_pct - lo_pct) * chart_w, 0.05)

        # Row background
        _rect(s4, 0.15, ty, W - 0.30, row_h2, bg_col)
        # Label
        _txt(s4, label, 0.22, ty + 0.08, 1.40, 0.20, size=7.5, color=LIGHT, bold=True)
        _txt(s4, f"{cs}{lo:,.1f}", 0.22, ty + 0.30, 1.40, 0.22, size=7, color=DIM)
        # Track
        _rect(s4, chart_lx, ty + 0.18, chart_w, bar_h, DIMMER)
        # Filled bar
        bc = bar_colors[ri]
        _rect(s4, bar_lx, ty + 0.18, bar_w, bar_h, bc)
        # Range labels on bar
        _txt(s4, f"{cs}{lo:,.0f}", bar_lx, ty + 0.50, 1.0, 0.20, size=6.5, color=bc)
        _txt(s4, f"{cs}{hi:,.0f}", bar_lx + bar_w - 0.9, ty + 0.50, 0.9, 0.20, size=6.5, color=bc, align=PP_ALIGN.RIGHT)

    # Current price vertical line
    cur_pct = (cur_price - g_min) / g_range
    cur_lx  = chart_lx + cur_pct * chart_w
    _rect(s4, cur_lx - 0.01, row_start_y, 0.03, len(rows) * (row_h2 + 0.04), WHITE)
    _txt(s4, f"▲ {cs}{cur_price:,.2f}", cur_lx - 0.35, row_start_y - 0.35, 0.80, 0.28,
         size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Valuation summary strip at bottom
    va = request.valuation_analysis or {}
    _rect(s4, 0.15, 3.95, W - 0.30, 0.90, PANEL)
    summary_vals = [
        ("DCF Value",    f"{cs}{va.get('dcf_price', '-')}"),
        ("Comps Value",  f"{cs}{va.get('comps_price', '-')}"),
        ("Blended",      f"{cs}{request.implied_price:.2f}"),
        ("Spread",       f"{request.upside:+.1f}%"),
        ("WACC",         f"{va.get('wacc', '-')}%"),
        ("Growth Assumed", f"{va.get('growth_assumed', '-')}%"),
    ]
    sv_w = (W - 0.60) / len(summary_vals)
    for si, (lbl, val) in enumerate(summary_vals):
        sx = 0.20 + si * sv_w
        _txt(s4, lbl.upper(), sx, 4.00, sv_w - 0.05, 0.22, size=6, color=DIM, align=PP_ALIGN.CENTER)
        v_c = GREEN if "+" in str(val) and si == 3 else (RED if "-" in str(val) and si == 3 else WHITE)
        _txt(s4, val, sx, 4.26, sv_w - 0.05, 0.32, size=11, bold=True, color=v_c, align=PP_ALIGN.CENTER)

    if not is_client:
        _rect(s4, 0.15, 5.00, W - 0.30, 0.45, HDR_BEAR)
        _txt(s4,
             f"STRESS SCENARIO  ·  DCF Range: {request.dcf_range}  ·  "
             f"WACC Stress: +{round((va.get('wacc') or 0) * 1.25, 1) if isinstance(va.get('wacc'), (int, float)) else 'N/A'}%  ·  "
             f"Bear Sensitivity: FCF −20%",
             0.30, 5.04, W - 0.60, 0.34, size=7.5, color=RED)

    _footer(s4)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — BULL / BEAR NARRATIVE  (Internal full detail | Client: highlights)
    # ════════════════════════════════════════════════════════════════════════════
    s5 = _slide()
    _bg(s5, BG)
    _rect(s5, 0, 0, W, 0.70, PANEL)
    _rect(s5, 0, 0, 0.10, H, ACCENT)

    theses = request.theses or {}

    if is_client:
        _txt(s5, f"{request.ticker}  ·  Investment Highlights",
             0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

        # Highlights panel
        swot = request.swot or {}
        highlights = (swot.get("Strengths") or []) + (swot.get("Opportunities") or [])
        _rect(s5, 0.20, 0.85, W - 0.40, 4.70, PANEL2)
        _rect(s5, 0.20, 0.85, 0.10, 4.70, GREEN)
        _txt(s5, "KEY INVESTMENT MERITS", 0.40, 0.92, 8.5, 0.28, size=8, bold=True, color=GREEN)
        for i, pt in enumerate(highlights[:8]):
            _txt(s5, f"→  {pt}", 0.40, 1.28 + i * 0.48, 9.1, 0.40, size=9.5, color=LIGHT)

        # Risk disclosure
        _rect(s5, 0.20, 5.70, W - 0.40, 0.70, HDR_BEAR)
        _txt(s5, "⚠  Risk Disclosure: Past performance is not indicative of future results. "
             "This material is for informational purposes only and does not constitute investment advice.",
             0.35, 5.78, W - 0.60, 0.52, size=7.5, color=YELLOW, wrap=True)

    else:
        _txt(s5, f"{request.ticker}  ·  Bull / Bear Thesis  (INTERNAL)",
             0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

        panel_w = (W - 0.60) / 2 - 0.10
        bull_x, bear_x = 0.20, 0.20 + panel_w + 0.20

        # BULL panel
        _rect(s5, bull_x, 0.82, panel_w, 4.30, PANEL2)
        _rect(s5, bull_x, 0.82, panel_w, 0.52, HDR_BULL)
        _rect(s5, bull_x, 0.82, 0.10, 4.30, GREEN)
        _txt(s5, "▲  BULL CASE", bull_x + 0.18, 0.88, panel_w - 0.22, 0.36,
             size=10, bold=True, color=GREEN)
        bull_text = theses.get("bull", "") or ""
        # Split into sentences for bullet rendering
        bull_sentences = [s.strip() for s in bull_text.replace(". ", ".|").split("|") if s.strip()]
        for i, sent in enumerate(bull_sentences[:6]):
            _txt(s5, f"•  {sent}", bull_x + 0.18, 1.44 + i * 0.55, panel_w - 0.28, 0.46,
                 size=8.5, color=LIGHT, wrap=True)

        # BEAR panel
        _rect(s5, bear_x, 0.82, panel_w, 4.30, PANEL2)
        _rect(s5, bear_x, 0.82, panel_w, 0.52, HDR_BEAR)
        _rect(s5, bear_x, 0.82, 0.10, 4.30, RED)
        _txt(s5, "▼  BEAR CASE", bear_x + 0.18, 0.88, panel_w - 0.22, 0.36,
             size=10, bold=True, color=RED)
        bear_text = theses.get("bear", "") or ""
        bear_sentences = [s.strip() for s in bear_text.replace(". ", ".|").split("|") if s.strip()]
        for i, sent in enumerate(bear_sentences[:6]):
            _txt(s5, f"•  {sent}", bear_x + 0.18, 1.44 + i * 0.55, panel_w - 0.28, 0.46,
                 size=8.5, color=LIGHT, wrap=True)

        # Stress scenarios strip
        _rect(s5, 0.20, 5.25, W - 0.40, 0.62, HDR_BEAR)
        _rect(s5, 0.20, 5.25, 0.08, 0.62, RED)
        va2 = request.valuation_analysis or {}
        _txt(s5,
             f"STRESS SCENARIOS (INTERNAL)   ·   "
             f"DCF Bear: {cs}{va2.get('dcf_price', 0) * 0.75:.2f}  "
             f"·  DCF Bull: {cs}{va2.get('dcf_price', 0) * 1.25:.2f}  "
             f"·  WACC Range: {request.wacc_components.get('wacc', '?')}% – "
             f"{round(float(request.wacc_components.get('wacc', 0) or 0) * 1.3, 1)}%  "
             f"·  FCF Sensitivity: −20% / +20%",
             0.35, 5.32, W - 0.55, 0.46, size=7.5, color=RED, wrap=True)

    _footer(s5)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — COMPARABLE ANALYSIS
    # ════════════════════════════════════════════════════════════════════════════
    s6 = _slide()
    _bg(s6, BG)
    _rect(s6, 0, 0, W, 0.70, PANEL)
    _rect(s6, 0, 0, 0.10, H, ACCENT)
    _txt(s6, f"{request.ticker}  ·  Comparable Company Analysis",
         0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

    peers = request.peers[:8]
    headers  = ["SYMBOL", "CAT", "MKT CAP", "P/E", "EV/EBITDA", "EV/SALES", "ROIC", "SIMILARITY"]
    col_ws   = [1.50, 0.55, 1.30, 0.85, 1.10, 1.10, 0.90, 1.35]
    x_starts = [0.20]
    for cw in col_ws[:-1]:
        x_starts.append(x_starts[-1] + cw + 0.02)

    row_h3  = 0.44
    hdr_ty  = 0.82
    # Header row
    _rect(s6, 0.18, hdr_ty, W - 0.36, row_h3, HDR_ROW)
    for hdr, cw, cx in zip(headers, col_ws, x_starts):
        _txt(s6, hdr, cx + 0.05, hdr_ty + 0.10, cw - 0.08, 0.24,
             size=7, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    for ri, peer in enumerate(peers):
        row_ty = hdr_ty + (ri + 1) * (row_h3 + 0.02)
        bg = STRIP_A if ri % 2 == 0 else STRIP_B
        _rect(s6, 0.18, row_ty, W - 0.36, row_h3, bg)

        # Category badge color
        cat = peer.category or "A"
        cat_col = ACCENT if cat == "A" else YELLOW

        vals = [
            (peer.symbol, WHITE, True),
            (cat, cat_col, True),
            (peer.mkt_cap, LIGHT, False),
            (f"{peer.pe:.1f}x" if peer.pe else "-", LIGHT, False),
            (f"{peer.ev_ebitda:.1f}x" if peer.ev_ebitda else "-", LIGHT, False),
            (f"{peer.ev_sales:.1f}x" if peer.ev_sales else "-", LIGHT, False),
            (peer.roic, LIGHT, False),
            (f"{peer.similarity:.0f}%", GREEN if peer.similarity >= 70 else (YELLOW if peer.similarity >= 40 else RED), True),
        ]
        for (val, vc, b), cw, cx in zip(vals, col_ws, x_starts):
            _txt(s6, str(val), cx + 0.05, row_ty + 0.10, cw - 0.08, 0.26,
                 size=8, bold=b, color=vc)

    # Category legend
    _txt(s6, "CAT A = Direct Sector & Industry Comparables   ·   CAT B = Sector Scale Benchmarks",
         0.20, H - 1.02, W - 0.40, 0.24, size=6.5, color=DIM, italic=True)
    _footer(s6)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — SWOT ANALYSIS  (Internal) | Risk Disclosures (Client)
    # ════════════════════════════════════════════════════════════════════════════
    s7 = _slide()
    _bg(s7, BG)
    _rect(s7, 0, 0, W, 0.70, PANEL)
    _rect(s7, 0, 0, 0.10, H, ACCENT)

    swot = request.swot or {}

    if is_client:
        _txt(s7, f"{request.ticker}  ·  Risk Disclosures",
             0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)
        risks = [
            "Market Risk: Equity markets can be volatile; share prices may decline significantly.",
            "Liquidity Risk: Ability to buy/sell at quoted prices is not guaranteed.",
            "Model Risk: Valuation models are based on assumptions that may not materialise.",
            "Regulatory Risk: Changes in laws or regulation may adversely affect the company.",
            "Macro Risk: Interest rate changes, inflation, and geopolitical events may impact returns.",
            "Currency Risk: International investments are subject to exchange rate fluctuations.",
            "Information Risk: Financial data is sourced from public disclosures which may contain errors.",
        ]
        for i, risk in enumerate(risks):
            _rect(s7, 0.20, 0.85 + i * 0.72, W - 0.40, 0.62, PANEL2 if i % 2 == 0 else PANEL)
            _rect(s7, 0.20, 0.85 + i * 0.72, 0.06, 0.62, RED)
            _txt(s7, f"• {risk}", 0.34, 0.90 + i * 0.72, W - 0.55, 0.48,
                 size=8.5, color=LIGHT, wrap=True)
    else:
        _txt(s7, f"{request.ticker}  ·  SWOT Analysis  (INTERNAL)",
             0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

        quad_w = (W - 0.60) / 2
        quad_h = (H - 1.20) / 2 - 0.08
        positions = {
            "Strengths":     (0.20,            0.82,  GREEN,  "▲ STRENGTHS"),
            "Weaknesses":    (0.20 + quad_w + 0.18, 0.82,  RED,    "▼ WEAKNESSES"),
            "Opportunities": (0.20,            0.82 + quad_h + 0.12, ACCENT, "◆ OPPORTUNITIES"),
            "Threats":       (0.20 + quad_w + 0.18, 0.82 + quad_h + 0.12, YELLOW, "⚡ THREATS"),
        }
        for section, (lx, ty, col, label) in positions.items():
            _rect(s7, lx, ty, quad_w, quad_h, PANEL2)
            _rect(s7, lx, ty, quad_w, 0.44, col)
            _txt(s7, label, lx + 0.12, ty + 0.08, quad_w - 0.18, 0.28,
                 size=8.5, bold=True, color=BG)
            items_list = swot.get(section) or ["N/A"]
            for ii, item in enumerate(items_list[:4]):
                _txt(s7, f"→  {item}", lx + 0.12, ty + 0.55 + ii * 0.48,
                     quad_w - 0.22, 0.42, size=8, color=LIGHT, wrap=True)

    _footer(s7)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — WACC & MODEL DETAIL  (Internal) | Closing (Client)
    # ════════════════════════════════════════════════════════════════════════════
    s8 = _slide()
    _bg(s8, BG)
    _rect(s8, 0, 0, W, 0.70, PANEL)
    _rect(s8, 0, 0, 0.10, H, ACCENT)

    if is_client:
        # Client closing slide
        _txt(s8, "Thank You", 0.30, 1.80, 9.0, 1.20, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s8, f"For questions regarding this analysis, please contact your coverage team.",
             0.30, 3.20, 9.0, 0.50, size=10, color=DIM, align=PP_ALIGN.CENTER)
        _txt(s8, "SAMPADA.ai  ·  Institutional Equity Research  ·  Educational Simulation Only",
             0.30, 4.00, 9.0, 0.34, size=8, color=DIMMER, align=PP_ALIGN.CENTER)
    else:
        _txt(s8, f"{request.ticker}  ·  WACC Bridge & Model Components  (INTERNAL)",
             0.25, 0.17, 9.0, 0.38, size=13, bold=True, color=WHITE)

        wc = request.wacc_components or {}
        _section_header(s8, "COST OF CAPITAL BRIDGE", 0.82)

        wacc_items = [
            ("Risk-Free Rate (10Y)",   f"{round(float(wc.get('rf') or 0) * 100, 2)}%"),
            ("Beta (Blume Adjusted)",  f"{round(float(wc.get('beta') or 1.0), 2)}x"),
            ("Equity Risk Premium",    f"{round(float(wc.get('erp') or 0.055) * 100, 2)}%"),
            ("WACC (Base Case)",       f"{wc.get('wacc', '-')}%"),
        ]
        for i, (lbl, val) in enumerate(wacc_items):
            _kv_box(s8, lbl, val, lx=0.20 + i * 2.35, ty=1.12, w=2.25, h=0.85,
                    val_color=ACCENT if "WACC" in lbl else WHITE)

        _section_header(s8, "DCF MODEL ASSUMPTIONS", 2.15)
        va3 = request.valuation_analysis or {}
        dcf_items = [
            ("Explicit Growth Period",  "10 Years"),
            ("Initial Growth Rate",     f"{va3.get('growth_assumed', '-')}%"),
            ("Terminal Growth Rate",    "2.5%"),
            ("FCF Blend",               "70% Perpetuity + 30% Exit"),
            ("Blend Weight",            "60% DCF / 40% Comps"),
        ]
        for i, (lbl, val) in enumerate(dcf_items):
            _kv_box(s8, lbl, val, lx=0.20 + i * 1.92, ty=2.42, w=1.82, h=0.75,
                    val_color=LIGHT)

        _section_header(s8, "PEER VALUATION MULTIPLES", 3.35)
        va_rows = [
            ("Metric Used",   dealtype_label := {"M&A": "P/E", "IPO": "EV/Sales", "LBO": "EV/EBITDA"}.get(request.deal_type, "P/E")),
            ("Deal Type",     request.deal_type),
            ("Peers (Cat A)", str(sum(1 for p in request.peers if (p.category or "A") == "A"))),
            ("Peers (Cat B)", str(sum(1 for p in request.peers if (p.category or "A") == "B"))),
        ]
        for i, (lbl, val) in enumerate(va_rows):
            _kv_box(s8, lbl, val, lx=0.20 + i * 2.35, ty=3.60, w=2.25, h=0.75)

        # Methodology note
        _rect(s8, 0.20, 4.55, W - 0.40, 0.90, PANEL)
        _txt(s8,
             "Methodology: Institutional DCF with dual-scenario WACC (Blume-adjusted beta, sector ERP). "
             "Comps via Similarity-Weighted Harmonic Mean across direct (Cat A) and scale (Cat B) peers. "
             "Final value: 60% intrinsic DCF + 40% relative comps. Terminal value: Gordon Growth + EBITDA Exit blend.",
             0.35, 4.62, W - 0.60, 0.75, size=7.5, color=DIM, wrap=True)

    _footer(s8)

    # ── Serialise & stream ─────────────────────────────────────────────────────
    output = BytesIO()
    prs.save(output)
    output.seek(0)

    mode_tag = "Client" if is_client else "Internal"
    fname = f"SAMPADA_{request.ticker}_{mode_tag}_{time.strftime('%Y%m%d')}.pptx"
    return StreamingResponse(output, headers={
        'Content-Disposition': f'attachment; filename="{fname}"',
        'Content-Type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    })
